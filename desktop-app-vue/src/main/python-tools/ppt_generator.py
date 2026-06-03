#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成工具
使用python-pptx库创建演示文稿
"""

import sys
import json
import os
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print(json.dumps({
        'success': False,
        'error': '未安装python-pptx，请运行: pip install python-pptx'
    }, ensure_ascii=False))
    sys.exit(1)

class PPTGenerator:
    """PPT生成器"""

    def __init__(self):
        self.prs = Presentation()
        # 设置幻灯片尺寸（16:9）
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def create_presentation(self, params):
        """
        创建PPT演示文稿

        参数:
        - title: 演示文稿标题
        - subtitle: 副标题
        - slides: 幻灯片列表 [{ type, title, content, layout }]
        - output_path: 输出路径
        - template: 模板类型 (business/education/creative)
        """
        title = params.get('title', '未命名演示')
        subtitle = params.get('subtitle', '')
        slides = params.get('slides', [])
        output_path = params.get('output_path')
        template = params.get('template', 'business')

        # 创建封面
        self._create_title_slide(title, subtitle, template)

        # 创建内容幻灯片
        for slide_data in slides:
            self._create_content_slide(slide_data, template)

        # 保存
        if not output_path:
            output_path = f"{title}.pptx"

        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)

        self.prs.save(output_path)

        return {
            'success': True,
            'output_path': output_path,
            'file_size': os.path.getsize(output_path),
            'slides_count': len(self.prs.slides),
            'created_at': datetime.now().isoformat()
        }

    def _create_title_slide(self, title, subtitle, template):
        """创建封面幻灯片"""
        slide_layout = self.prs.slide_layouts[0]  # 标题幻灯片布局
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle if subtitle else datetime.now().strftime('%Y年%m月%d日')

        # 根据模板设置样式
        if template == 'business':
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        elif template == 'education':
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(46, 125, 50)
            title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        elif template == 'creative':
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(156, 39, 176)
            title_shape.text_frame.paragraphs[0].font.size = Pt(48)

    def _create_content_slide(self, slide_data, template):
        """创建内容幻灯片"""
        slide_type = slide_data.get('type', 'title_content')
        slide_title = slide_data.get('title', '')
        content = slide_data.get('content', [])

        if slide_type == 'title_only':
            layout = self.prs.slide_layouts[5]  # 仅标题
        elif slide_type == 'two_content':
            layout = self.prs.slide_layouts[3]  # 两栏内容
        else:
            layout = self.prs.slide_layouts[1]  # 标题和内容

        slide = self.prs.slides.add_slide(layout)

        # 设置标题
        if slide.shapes.title:
            slide.shapes.title.text = slide_title

        # 添加内容
        if slide_type == 'title_content':
            self._add_bullet_points(slide, content)
        elif slide_type == 'two_content':
            self._add_two_column_content(slide, content)
        elif slide_type == 'image':
            self._add_image_slide(slide, content)

    def _add_bullet_points(self, slide, content):
        """添加项目符号列表"""
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            if isinstance(content, list):
                for item in content:
                    p = tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0
                    p.font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = str(content)
                p.font.size = Pt(18)

    def _add_two_column_content(self, slide, content):
        """添加两栏内容"""
        left_content = content.get('left', [])
        right_content = content.get('right', [])

        # 左栏位置
        left_left = Inches(0.5)
        left_top = Inches(1.5)
        left_width = Inches(4.5)
        left_height = Inches(3.5)

        # 右栏位置
        right_left = Inches(5.2)
        right_top = Inches(1.5)
        right_width = Inches(4.5)
        right_height = Inches(3.5)

        # 添加左栏文本框
        left_box = slide.shapes.add_textbox(left_left, left_top, left_width, left_height)
        left_tf = left_box.text_frame
        left_tf.word_wrap = True

        for i, item in enumerate(left_content):
            p = left_tf.add_paragraph() if i > 0 else left_tf.paragraphs[0]
            p.text = str(item)
            p.font.size = Pt(16)
            p.level = 0

        # 添加右栏文本框
        right_box = slide.shapes.add_textbox(right_left, right_top, right_width, right_height)
        right_tf = right_box.text_frame
        right_tf.word_wrap = True

        for i, item in enumerate(right_content):
            p = right_tf.add_paragraph() if i > 0 else right_tf.paragraphs[0]
            p.text = str(item)
            p.font.size = Pt(16)
            p.level = 0

    def _add_image_slide(self, slide, content):
        """添加图片幻灯片"""
        image_path = content.get('image_path', content.get('path', ''))
        caption = content.get('caption', '')

        if not image_path or not os.path.exists(image_path):
            # 如果图片不存在，添加占位文本
            placeholder = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            tf = placeholder.text_frame
            p = tf.paragraphs[0]
            p.text = f"[图片: {image_path}]"
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER
            return

        # 计算图片位置（居中）
        img_left = Inches(1.5)
        img_top = Inches(1.3)
        img_width = Inches(7)
        img_height = Inches(3.5)

        # 添加图片
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

        # 添加图片说明
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
            tf = caption_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER

def main():
    """主函数"""
    try:
        if len(sys.argv) < 2:
            raise ValueError('缺少参数')

        args = json.loads(sys.argv[1])

        generator = PPTGenerator()

        operation = args.get('operation', 'create')

        if operation == 'create':
            result = generator.create_presentation(args)
        else:
            raise ValueError(f"不支持的操作: {operation}")

        print(json.dumps(result, ensure_ascii=False, indent=2))

    except Exception as e:
        error_result = {
            'success': False,
            'error': str(e),
            'type': type(e).__name__
        }
        print(json.dumps(error_result, ensure_ascii=False))
        sys.exit(1)

if __name__ == '__main__':
    main()
